from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import tempfile
import os
import math
from typing import List, Dict, Optional

# ── Palette definitions ────────────────────────────────────────────────────────
PALETTES = {
    "academic": {
        "primary":          RGBColor(0, 51, 102),
        "secondary":        RGBColor(255, 255, 255),
        "accent":           RGBColor(212, 175, 55),
        "dark":             RGBColor(26, 26, 26),
        "light_bg":         RGBColor(235, 242, 255),
        "def_box":          RGBColor(210, 228, 255),
        "err_box":          RGBColor(255, 220, 220),
        "mid_gray":         RGBColor(150, 150, 150),
        "col_right_header": RGBColor(0, 110, 60),
        "col_right_box":    RGBColor(220, 245, 225),
        "col_right_txt":    RGBColor(26, 26, 26),
    },
    "professional": {
        "primary":          RGBColor(0, 80, 100),
        "secondary":        RGBColor(255, 255, 255),
        "accent":           RGBColor(230, 100, 50),
        "dark":             RGBColor(51, 51, 51),
        "light_bg":         RGBColor(230, 250, 250),
        "def_box":          RGBColor(200, 240, 240),
        "err_box":          RGBColor(255, 220, 220),
        "mid_gray":         RGBColor(150, 150, 150),
        "col_right_header": RGBColor(180, 70, 20),
        "col_right_box":    RGBColor(255, 235, 210),
        "col_right_txt":    RGBColor(51, 51, 51),
    },
    "clean": {
        "primary":          RGBColor(75, 0, 130),
        "secondary":        RGBColor(255, 255, 255),
        "accent":           RGBColor(0, 180, 130),
        "dark":             RGBColor(26, 26, 26),
        "light_bg":         RGBColor(240, 232, 255),
        "def_box":          RGBColor(225, 210, 255),
        "err_box":          RGBColor(255, 220, 240),
        "mid_gray":         RGBColor(150, 150, 150),
        "col_right_header": RGBColor(0, 150, 110),
        "col_right_box":    RGBColor(210, 248, 238),
        "col_right_txt":    RGBColor(26, 26, 26),
    },
    "warm": {
        "primary":          RGBColor(139, 0, 0),
        "secondary":        RGBColor(255, 255, 255),
        "accent":           RGBColor(255, 200, 0),
        "dark":             RGBColor(26, 26, 26),
        "light_bg":         RGBColor(255, 240, 235),
        "def_box":          RGBColor(255, 220, 210),
        "err_box":          RGBColor(255, 210, 210),
        "mid_gray":         RGBColor(150, 150, 150),
        "col_right_header": RGBColor(180, 100, 0),
        "col_right_box":    RGBColor(255, 238, 200),
        "col_right_txt":    RGBColor(26, 26, 26),
    },
}

SUBJECT_PALETTE = {
    "research": "academic", "methodology": "academic",
    "science": "professional", "biology": "professional",
    "chemistry": "professional", "physics": "professional",
    "math": "clean", "statistics": "clean",
    "computer": "clean", "technology": "clean",
    "history": "warm", "filipino": "warm",
    "literature": "warm", "economics": "academic",
    "business": "academic",
}


def get_palette(subject: str) -> dict:
    for key, name in SUBJECT_PALETTE.items():
        if key in subject.lower():
            return PALETTES[name]
    return PALETTES["academic"]


# ── Slide dimensions ───────────────────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)
HEADER_H   = Inches(1.1)
ACCENT_H   = Inches(0.08)
CONTENT_TOP = HEADER_H + ACCENT_H + Inches(0.12)
MARGIN_L   = Inches(0.4)
MARGIN_R   = Inches(0.4)
CONTENT_W  = W - MARGIN_L - MARGIN_R
FOOTER_Y   = Inches(7.1)


# ── Font sizing helpers ────────────────────────────────────────────────────────
def fit_font(items: list, box_w_in: float, box_h_in: float,
             max_pt: float = 28, min_pt: float = 14,
             spacing_in: float = 0.18) -> float:
    def total_height(pt):
        chars_per_line = max(1, box_w_in / (pt * 0.0064))
        line_h = pt * 0.0175
        h = 0.0
        for item in items:
            lines = math.ceil(max(1, len(item) / chars_per_line))
            h += lines * line_h
        h += spacing_in * max(0, len(items) - 1)
        return h

    pt = max_pt
    while pt > min_pt:
        if total_height(pt) <= box_h_in:
            return round(pt, 1)
        pt -= 0.5
    return min_pt


def fit_title_font(text: str, box_w_in: float, box_h_in: float,
                   max_pt: float = 48, min_pt: float = 22) -> float:
    def h(pt):
        chars_per_line = max(1, box_w_in / (pt * 0.0064))
        line_h = pt * 0.0175
        lines = math.ceil(max(1, len(text) / chars_per_line))
        return lines * line_h
    pt = max_pt
    while pt > min_pt:
        if h(pt) <= box_h_in:
            return round(pt, 1)
        pt -= 0.5
    return min_pt


def max_pt_for_n(n: int) -> float:
    if n <= 1: return 28.0
    if n == 2: return 26.0
    return 24.0


# ── Shape helpers ──────────────────────────────────────────────────────────────
def add_rect(slide, left, top, width, height,
             fill=None, line=None, line_pt=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_pt:
            shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_tb(slide, left, top, width, height,
           text="", size=18, bold=False, color=None,
           align=PP_ALIGN.LEFT, wrap=True, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    if color:
        run.font.color.rgb = color
    return tb, tf


def vcenter_tb(tb):
    from pptx.oxml.ns import qn
    txBody = tb.text_frame._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("anchor", "ctr")


def add_para(tf, text, size=18, bold=False, color=None,
             font="Calibri", space_before=0):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    if color:
        run.font.color.rgb = color
    return p


def write_notes(pptx_slide, notes_text: str):
    if not notes_text:
        return
    try:
        notes_slide = pptx_slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text
    except Exception:
        pass


# ── Icon mapper ────────────────────────────────────────────────────────────────
ICON_MAP = [
    (["growth", "increase", "improve", "advantage", "benefit",
      "success", "achieve", "develop", "expand", "opportunity", "progress"], "📈"),
    (["challenge", "problem", "issue", "difficulty", "barrier",
      "obstacle", "risk", "concern", "weakness", "disadvantage"], "⚠️"),
    (["data", "research", "study", "analysis", "statistic", "survey",
      "finding", "result", "evidence", "method", "framework", "theory"], "🔬"),
    (["student", "teacher", "community", "family", "people",
      "society", "culture", "social", "human", "filipino"], "👥"),
    (["education", "school", "university", "college", "learning",
      "curriculum", "academic", "classroom", "ched", "deped"], "🎓"),
    (["technology", "digital", "software", "internet", "computer",
      "online", "virtual", "ai", "automation", "platform"], "💻"),
    (["economy", "economic", "gdp", "inflation", "finance",
      "investment", "income", "poverty", "budget", "trade"], "💰"),
    (["policy", "government", "law", "regulation", "legislation",
      "ched", "deped", "memorandum", "act", "republic"], "📋"),
    (["ethics", "ethical", "moral", "value", "principle",
      "justice", "integrity", "respect", "consent", "privacy"], "⚖️"),
    (["definition", "concept", "term", "meaning", "refers", "defined", "known as"], "📖"),
    (["process", "step", "phase", "stage", "procedure", "strategy",
      "implement", "conduct", "collect", "analyze"], "🔄"),
    (["history", "historical", "origin", "founded", "established",
      "century", "decade", "period", "evolution"], "📅"),
    (["philippine", "pilipino", "manila", "mindanao", "luzon",
      "visayas", "bayanihan", "ofw"], "🇵🇭"),
    ([], "▶️"),
]


def get_icon(text: str) -> str:
    lower = text.lower()
    for keywords, icon in ICON_MAP:
        if not keywords:
            return icon
        if any(kw in lower for kw in keywords):
            return icon
    return "▶️"


def add_icon(text: str) -> str:
    if not text:
        return text
    if ord(text[0]) > 127:
        return text
    return f"{get_icon(text)}  {text}"


# ── Common header ──────────────────────────────────────────────────────────────
def build_header(slide, heading: str, pal: dict):
    add_rect(slide, 0, 0, W, HEADER_H, fill=pal["primary"])
    add_rect(slide, 0, HEADER_H, W, ACCENT_H, fill=pal["accent"])
    add_tb(slide, Inches(0.5), Inches(0.18), Inches(12.333), Inches(0.8),
           text=heading, size=28, bold=True,
           color=pal["secondary"], align=PP_ALIGN.LEFT, font="Cambria")


# ── Title slide ────────────────────────────────────────────────────────────────
def build_title_slide(slide, data: dict, title: str, subject: str, level: str, pal: dict):
    from datetime import datetime
    add_rect(slide, 0, 0, W, H, fill=pal["primary"])
    add_rect(slide, 0, 0, W, Inches(0.18), fill=pal["accent"])
    add_rect(slide, 0, Inches(7.32), W, Inches(0.18), fill=pal["accent"])

    BOX_L = Inches(1.0)
    BOX_T = Inches(1.3)
    BOX_W = Inches(11.333)
    BOX_H = Inches(5.0)

    box = add_rect(slide, BOX_L, BOX_T, BOX_W, BOX_H, fill=RGBColor(255, 255, 255))
    box.line.color.rgb = pal["accent"]
    box.line.width = Pt(3)

    title_pt = fit_title_font(title, 10.9, 2.4, max_pt=44, min_pt=22)
    add_tb(slide, BOX_L + Inches(0.2), BOX_T + Inches(0.2),
           Inches(10.9), Inches(2.4),
           text=title, size=title_pt, bold=True, color=pal["primary"],
           align=PP_ALIGN.CENTER, font="Cambria")

    DIV_Y = BOX_T + Inches(2.65)
    add_rect(slide, BOX_L + Inches(1.5), DIV_Y,
             Inches(8.333), Inches(0.05), fill=pal["accent"])

    add_tb(slide, BOX_L + Inches(0.2), DIV_Y + Inches(0.15),
           Inches(10.9), Inches(0.6),
           text=f"{subject}  •  {level} Level",
           size=20, color=RGBColor(80, 80, 80),
           align=PP_ALIGN.CENTER, font="Calibri")

    PILL_Y = DIV_Y + Inches(0.9)
    add_rect(slide, Inches(5.0), PILL_Y, Inches(3.333), Inches(0.5),
             fill=pal["accent"])
    add_tb(slide, Inches(5.0), PILL_Y + Inches(0.04),
           Inches(3.333), Inches(0.44),
           text=datetime.now().strftime("%B %d, %Y"),
           size=14, bold=True, color=pal["primary"],
           align=PP_ALIGN.CENTER, font="Calibri")

    add_tb(slide, Inches(0.5), Inches(7.05), Inches(12.333), Inches(0.3),
           text="🎓 pptPro — AI-Powered Academic Presentations for Filipino Students",
           size=11, color=RGBColor(200, 200, 200), align=PP_ALIGN.CENTER)


# ── Section divider ────────────────────────────────────────────────────────────
def build_section_slide(slide, number: int, title: str, description: str, pal: dict):
    add_rect(slide, 0, 0, W, H, fill=pal["primary"])
    add_rect(slide, 0, 0, W, Inches(0.1), fill=pal["accent"])
    add_rect(slide, 0, Inches(7.4), W, Inches(0.1), fill=pal["accent"])

    circle = slide.shapes.add_shape(9, Inches(1.2), Inches(1.8), Inches(3.2), Inches(3.2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = pal["accent"]
    circle.line.fill.background()

    add_tb(slide, Inches(1.2), Inches(2.2), Inches(3.2), Inches(2.0),
           text=f"0{number}", size=76, bold=True,
           color=pal["secondary"], align=PP_ALIGN.CENTER, font="Cambria")

    add_tb(slide, Inches(5.0), Inches(2.2), Inches(7.8), Inches(1.5),
           text=title, size=40, bold=True,
           color=pal["secondary"], align=PP_ALIGN.LEFT, font="Cambria")

    add_tb(slide, Inches(5.0), Inches(3.9), Inches(7.8), Inches(1.0),
           text=description, size=20,
           color=RGBColor(210, 210, 210), align=PP_ALIGN.LEFT, font="Calibri")


# ── Bullets slide ──────────────────────────────────────────────────────────────
def build_bullets_slide(slide, data: dict, pal: dict):
    add_rect(slide, 0, 0, W, H, fill=RGBColor(255, 255, 255))
    build_header(slide, data.get("heading", "Content"), pal)

    items = data.get("content", [])[:3]
    if not items:
        return

    n = len(items)
    layout_n = min(data.get("_original_n", n), 3)
    gap_in = 0.18
    box_h_in = (5.65 - gap_in * (layout_n - 1)) / layout_n
    pad = 0.14
    txt_w_in = 11.8

    for i, item in enumerate(items):
        y_in = (CONTENT_TOP / 914400) + i * (box_h_in + gap_in)
        y = Inches(y_in)
        bh = Inches(box_h_in)
        font_pt = fit_font([item],
                           box_w_in=txt_w_in - 0.3,
                           box_h_in=box_h_in - pad * 2,
                           max_pt=max_pt_for_n(layout_n), min_pt=16)
        add_rect(slide, Inches(0.35), y + Inches(pad),
                 Inches(0.12), bh - Inches(pad * 2),
                 fill=pal["accent"])
        tb, _ = add_tb(slide, Inches(0.65), y + Inches(0.08),
                       Inches(txt_w_in), bh - Inches(0.16),
                       text=add_icon(item), size=font_pt,
                       color=pal["dark"], wrap=True, font="Calibri")
        vcenter_tb(tb)


# ── Numbered slide ─────────────────────────────────────────────────────────────
def build_numbered_slide(slide, data: dict, pal: dict):
    add_rect(slide, 0, 0, W, H, fill=RGBColor(255, 255, 255))
    build_header(slide, data.get("heading", "Key Points"), pal)

    items = data.get("content", [])[:3]
    if not items:
        return

    n = len(items)
    layout_n = min(data.get("_original_n", n), 3)
    gap_in = 0.18
    box_h_in = (5.65 - gap_in * (layout_n - 1)) / layout_n
    BADGE_W = 0.75
    TEXT_W = 10.9
    pad = 0.14

    for i, item in enumerate(items):
        y_in = (CONTENT_TOP / 914400) + i * (box_h_in + gap_in)
        y = Inches(y_in)
        bh = Inches(box_h_in)
        font_pt = fit_font([item],
                           box_w_in=TEXT_W - 0.3,
                           box_h_in=box_h_in - pad * 2,
                           max_pt=max_pt_for_n(layout_n), min_pt=16)

        badge = slide.shapes.add_shape(9, Inches(0.35), y, Inches(BADGE_W), bh)
        badge.fill.solid()
        badge.fill.fore_color.rgb = pal["primary"]
        badge.line.fill.background()
        add_tb(slide, Inches(0.35), y + bh * 0.3, Inches(BADGE_W), bh * 0.45,
               text=str(i + 1), size=min(30, int(box_h_in * 18)),
               bold=True, color=pal["secondary"],
               align=PP_ALIGN.CENTER, font="Cambria")

        add_rect(slide, Inches(1.25), y, Inches(TEXT_W), bh,
                 fill=pal["light_bg"], line=pal["primary"], line_pt=0.5)
        tb, _ = add_tb(slide, Inches(1.45), y + Inches(pad),
                       Inches(TEXT_W - 0.3), bh - Inches(pad * 2),
                       text=add_icon(item), size=font_pt,
                       color=pal["dark"], wrap=True, font="Calibri")
        vcenter_tb(tb)


# ── Definition (key terms) slide ───────────────────────────────────────────────
def build_definition_slide(slide, data: dict, pal: dict):
    add_rect(slide, 0, 0, W, H, fill=RGBColor(255, 255, 255))
    build_header(slide, "📖 " + data.get("heading", "Key Terms"), pal)

    items = data.get("content", [])[:3]
    if not items:
        return

    layout_n = min(data.get("_original_n", len(items)), 3)
    gap_in = 0.18
    box_h_in = (5.65 - gap_in * (layout_n - 1)) / layout_n
    pad = 0.14

    for i, item in enumerate(items):
        y_in = (CONTENT_TOP / 914400) + i * (box_h_in + gap_in)
        y = Inches(y_in)
        bh = Inches(box_h_in)
        font_pt = fit_font([item],
                           box_w_in=11.8,
                           box_h_in=box_h_in - pad * 2,
                           max_pt=max_pt_for_n(layout_n), min_pt=16)
        add_rect(slide, Inches(0.4), y, Inches(12.53), bh,
                 fill=pal["def_box"], line=pal["primary"], line_pt=1.2)
        tb, _ = add_tb(slide, Inches(0.65), y + Inches(pad),
                       Inches(12.0), bh - Inches(pad * 2),
                       text=add_icon(item), size=font_pt,
                       color=pal["dark"], wrap=True, font="Calibri")
        vcenter_tb(tb)


# ── Two-column slide ───────────────────────────────────────────────────────────
def build_two_column_slide(slide, data: dict, pal: dict,
                            left_label="Column A", right_label="Column B"):
    add_rect(slide, 0, 0, W, H, fill=RGBColor(255, 255, 255))
    build_header(slide, data.get("heading", "Comparison"), pal)

    ll = data.get("left_label") or left_label
    rl = data.get("right_label") or right_label

    content = data.get("content", [])
    mid = len(content) // 2
    left_items = content[:mid][:3]
    right_items = content[mid:][:3]

    PILL_H = Inches(0.55)
    PILL_TOP = CONTENT_TOP
    COL_TOP = PILL_TOP + PILL_H + Inches(0.08)
    COL_BOT = Inches(7.0)
    COL_H = COL_BOT - COL_TOP
    COL_W = Inches(6.0)
    COL_W_IN = 6.0

    right_hdr = pal.get("col_right_header", pal["accent"])
    right_box = pal.get("col_right_box", RGBColor(220, 245, 225))
    right_txt = pal.get("col_right_txt", pal["dark"])

    def draw_col(x_in, items, hdr_color, box_color, txt_color, label):
        if not items:
            return
        pill = slide.shapes.add_shape(5, Inches(x_in), PILL_TOP, COL_W, PILL_H)
        pill.fill.solid()
        pill.fill.fore_color.rgb = hdr_color
        pill.line.fill.background()
        add_tb(slide, Inches(x_in), PILL_TOP + Inches(0.04),
               COL_W, PILL_H,
               text=label, size=18, bold=True,
               color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER, font="Cambria")

        add_rect(slide, Inches(x_in), COL_TOP, COL_W, COL_H,
                 fill=box_color, line=hdr_color, line_pt=1)

        n_items = len(items)
        gap_in = 0.14
        slot_h = (COL_H / 914400 - gap_in * max(0, n_items - 1)) / n_items if n_items > 0 else COL_H / 914400
        pad_in = 0.1

        for i, item in enumerate(items):
            sy_in = COL_TOP / 914400 + i * (slot_h + gap_in)
            item_pt = fit_font([f"✓  {item}"],
                               box_w_in=COL_W_IN - 0.3,
                               box_h_in=slot_h - pad_in * 2,
                               max_pt=max_pt_for_n(n_items), min_pt=16)
            add_tb(slide, Inches(x_in + 0.18), Inches(sy_in + pad_in),
                   Inches(COL_W_IN - 0.3), Inches(slot_h - pad_in * 2),
                   text=f"✓  {item}", size=item_pt,
                   color=txt_color, wrap=True, font="Calibri")

    draw_col(0.4, left_items, pal["primary"], pal["def_box"], pal["dark"], ll)
    draw_col(6.93, right_items, right_hdr, right_box, right_txt, rl)


# ── Misconceptions slide ───────────────────────────────────────────────────────
def build_misconceptions_slide(slide, data: dict, pal: dict):
    add_rect(slide, 0, 0, W, H, fill=RGBColor(255, 255, 255))
    build_header(slide, "⚠️ " + data.get("heading", "Common Misconceptions"), pal)

    content = data.get("content", [])
    pairs = [(content[i], content[i + 1])
             for i in range(0, len(content) - 1, 2)]
    if len(content) % 2 == 1:
        pairs.append((content[-1], "The correct understanding varies by context."))
    pairs = pairs[:3]

    if not pairs:
        return

    err_fill = pal.get("err_box", RGBColor(255, 220, 220))
    err_line = RGBColor(180, 40, 40)
    err_txt = RGBColor(130, 0, 0)
    ok_fill = pal.get("col_right_box", RGBColor(210, 245, 225))
    ok_line = pal.get("col_right_header", RGBColor(0, 120, 70))
    ok_txt = pal.get("col_right_txt", RGBColor(0, 80, 40))

    n = len(pairs)
    gap_in = 0.18
    row_h_in = (5.65 - gap_in * max(0, n - 1)) / n
    pad = 0.12

    all_texts = [p[0] for p in pairs] + [p[1] for p in pairs]
    font_pt = fit_font(all_texts, box_w_in=5.6, box_h_in=row_h_in - pad * 2,
                       max_pt=22, min_pt=14)

    for i, (mistake, correction) in enumerate(pairs):
        y_in = (CONTENT_TOP / 914400) + i * (row_h_in + gap_in)
        y = Inches(y_in)
        rh = Inches(row_h_in)
        add_rect(slide, Inches(0.4), y, Inches(6.0), rh,
                 fill=err_fill, line=err_line, line_pt=1.5)
        add_tb(slide, Inches(0.6), y + Inches(pad), Inches(5.7),
               rh - Inches(pad * 2),
               text=f"✗  {mistake}", size=font_pt, color=err_txt, wrap=True)
        add_rect(slide, Inches(6.93), y, Inches(6.0), rh,
                 fill=ok_fill, line=ok_line, line_pt=1.5)
        add_tb(slide, Inches(7.13), y + Inches(pad), Inches(5.7),
               rh - Inches(pad * 2),
               text=f"✓  {correction}", size=font_pt, color=ok_txt, wrap=True)


# ── Discussion slide ───────────────────────────────────────────────────────────
def build_discussion_slide(slide, data: dict, pal: dict):
    add_rect(slide, 0, 0, W, H, fill=RGBColor(255, 255, 255))
    build_header(slide, "💬 " + data.get("heading", "Discussion Questions"), pal)

    items = data.get("content", [])[:3]
    if not items:
        return

    layout_n = min(data.get("_original_n", len(items)), 3)
    gap_in = 0.18
    box_h_in = (5.65 - gap_in * (layout_n - 1)) / layout_n
    pad = 0.14

    for i, text in enumerate(items):
        formatted = f"Q{i + 1}.  {text}"
        y_in = (CONTENT_TOP / 914400) + i * (box_h_in + gap_in)
        y = Inches(y_in)
        bh = Inches(box_h_in)
        font_pt = fit_font([formatted],
                           box_w_in=11.8,
                           box_h_in=box_h_in - pad * 2,
                           max_pt=max_pt_for_n(layout_n), min_pt=16)
        add_rect(slide, Inches(0.4), y, Inches(12.53), bh,
                 fill=pal["light_bg"], line=pal["primary"], line_pt=1)
        add_tb(slide, Inches(0.65), y + Inches(pad),
               Inches(12.0), bh - Inches(pad * 2),
               text=formatted, size=font_pt,
               color=pal["primary"], wrap=True, font="Cambria")


# ── References slide ───────────────────────────────────────────────────────────
def build_references_slide(slide, data: dict, pal: dict):
    add_rect(slide, 0, 0, W, H, fill=RGBColor(255, 255, 255))
    build_header(slide, "📚 References (APA 7th Edition)", pal)

    content = data.get("content", [])
    if not content:
        return

    font_pt = fit_font(content, box_w_in=11.8, box_h_in=5.5,
                       max_pt=18, min_pt=12, spacing_in=0.14)

    _, tf = add_tb(slide, Inches(0.5), CONTENT_TOP + Inches(0.1),
                   Inches(12.333), Inches(5.5),
                   text="", size=font_pt, color=pal["dark"], wrap=True)

    for i, ref in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(10)
            p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = ref
        run.font.size = Pt(font_pt)
        run.font.name = "Calibri"
        run.font.color.rgb = pal["dark"]


# ── Conclusion slide ───────────────────────────────────────────────────────────
def build_conclusion_slide(slide, data: dict, pal: dict, title: str):
    add_rect(slide, 0, 0, W, H, fill=pal["primary"])
    add_rect(slide, 0, 0, W, Inches(0.1), fill=pal["accent"])
    add_rect(slide, 0, Inches(7.4), W, Inches(0.1), fill=pal["accent"])

    add_tb(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(1.1),
           text="🎯 " + data.get("heading", "Conclusion"),
           size=34, bold=True, color=pal["accent"],
           align=PP_ALIGN.LEFT, font="Cambria")

    items = data.get("content", [])[:3]
    if not items:
        return

    layout_n = min(data.get("_original_n", len(items)), 3)
    gap_in = 0.18
    top_y = 2.1
    bot_y = 6.7
    box_h_in = (bot_y - top_y - gap_in * max(0, layout_n - 1)) / layout_n
    pad = 0.14

    for i, item in enumerate(items):
        y_in = top_y + i * (box_h_in + gap_in)
        y = Inches(y_in)
        bh = Inches(box_h_in)
        font_pt = fit_font([item],
                           box_w_in=11.5,
                           box_h_in=box_h_in - pad * 2,
                           max_pt=max_pt_for_n(layout_n), min_pt=16)
        add_rect(slide, Inches(0.65), y + Inches(pad),
                 Inches(0.12), bh - Inches(pad * 2), fill=pal["accent"])
        add_tb(slide, Inches(0.95), y + Inches(pad),
               Inches(11.8), bh - Inches(pad * 2),
               text=add_icon(item), size=font_pt,
               color=RGBColor(255, 255, 255), wrap=True, font="Calibri")

    add_rect(slide, 0, Inches(6.85), W, Inches(0.55), fill=RGBColor(20, 40, 80))
    add_tb(slide, Inches(0.5), Inches(6.87), Inches(12.333), Inches(0.44),
           text="Salamat! 🙏   Thank you!   •   pptPro — Mabuhay ang Filipino Students!",
           size=14, color=pal["accent"], align=PP_ALIGN.CENTER)


# ── Pagination ─────────────────────────────────────────────────────────────────
PAGINATABLE = {
    "bullets", "content", "intro", "background", "concept1", "concept2",
    "concept3", "deepdive", "types", "examples", "casestudy", "trends",
    "challenges", "recommendations", "takeaways", "keyterms",
    "discussion", "misconceptions", "conclusion",
}

SECTION_BEFORE = {
    "intro":      (1, "Introduction",         "Understanding the topic"),
    "keyterms":   (2, "Key Concepts",          "Core principles and definitions"),
    "concept1":   (3, "Core Concepts",         "Deep dive into main ideas"),
    "types":      (4, "Types & Applications",  "Classifications and use cases"),
    "examples":   (5, "Practical Application", "Real-world Philippine examples"),
    "discussion": (6, "Class Discussion",      "Think critically together"),
    "takeaways":  (7, "Key Takeaways",         "Summary and conclusions"),
}


def paginate(slide_data: dict, page_size: int = 3) -> List[dict]:
    stype = slide_data.get("slide_type", "content")
    content = slide_data.get("content", [])

    if stype not in PAGINATABLE or len(content) <= page_size:
        return [slide_data]

    original_n = len(content)
    pages = []
    for i in range(0, len(content), page_size):
        chunk = content[i: i + page_size]
        page = {**slide_data, "content": chunk, "_original_n": original_n}
        if i > 0:
            page = {**page, "_is_continuation": True}
        pages.append(page)
    return pages


# ── Main builder ───────────────────────────────────────────────────────────────
def build_presentation(slides: List[Dict], title: str, subject: str,
                        level: str, language: str) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    pal = get_palette(subject)

    # Expand slides with pagination
    expanded: List[Dict] = []
    for sd in slides:
        expanded.extend(paginate(sd))

    for slide_data in expanded:
        stype = slide_data.get("slide_type", "content")

        # Add section divider if needed
        if stype in SECTION_BEFORE and not slide_data.get("_is_continuation"):
            sec = SECTION_BEFORE[stype]
            s = prs.slides.add_slide(blank)
            build_section_slide(s, sec[0], sec[1], sec[2], pal)

        # Add content slide
        s = prs.slides.add_slide(blank)

        try:
            if stype == "title":
                build_title_slide(s, slide_data, title, subject, level, pal)
            elif stype == "toc":
                build_numbered_slide(s, slide_data, pal)
            elif stype == "keyterms":
                build_definition_slide(s, slide_data, pal)
            elif stype == "comparison":
                ll = slide_data.get("left_label", "📊 Aspect A")
                rl = slide_data.get("right_label", "📊 Aspect B")
                build_two_column_slide(s, slide_data, pal, ll, rl)
            elif stype == "proscons":
                build_two_column_slide(s, slide_data, pal,
                                        "✅ Advantages", "❌ Disadvantages")
            elif stype == "misconceptions":
                build_misconceptions_slide(s, slide_data, pal)
            elif stype == "discussion":
                build_discussion_slide(s, slide_data, pal)
            elif stype == "references":
                build_references_slide(s, slide_data, pal)
            elif stype == "conclusion":
                build_conclusion_slide(s, slide_data, pal, title)
            elif stype in ("intro", "takeaways", "recommendations"):
                build_numbered_slide(s, slide_data, pal)
            else:
                build_bullets_slide(s, slide_data, pal)
        except Exception as e:
            print(f"⚠️ Slide build error for type '{stype}': {e}")
            # Render a safe fallback slide
            try:
                build_bullets_slide(s, slide_data, pal)
            except Exception:
                pass

        # Write speaker notes
        notes_text = slide_data.get("notes", "")
        if notes_text:
            write_notes(s, notes_text)

    # Save to temp file
    tmp = tempfile.NamedTemporaryFile(
    delete=False, suffix=".pptx", dir=tempfile.gettempdir()
)
    tmp_path = tmp.name
    tmp.close()

    prs.save(tmp_path)

    file_size = os.path.getsize(tmp_path)
    if file_size < 1000:
        raise RuntimeError(
            f"Generated PPTX is suspiciously small ({file_size} bytes). Please try again."
        )

    print(f"✅ PPTX saved: {tmp_path} ({file_size:,} bytes)")
    return tmp_path
